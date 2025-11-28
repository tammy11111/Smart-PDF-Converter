import streamlit as st
import cv2
import numpy as np
import pytesseract
from pytesseract import Output
from pdf2image import convert_from_bytes
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import io
import os
import re

# --- 網頁設定 ---
st.set_page_config(page_title="PDF 轉 PPT (段落增強版)", layout="wide")

st.title("📄 PDF 轉 PPT：段落增強 + 智慧過濾")
st.markdown("""
**本次修正：**
1. **段落敏感度提升**：即使只有一行，只要字數夠多或寬度夠寬，就會被視為內文拆解。
2. **繞圖排版支援**：主要內文即使緊貼圖片，也會被拆解，不會被誤判為圖說。
3. **小字保護**：只有「短小且緊貼圖片」的文字才會保留在背景。
""")

# --- 參數設定 ---
OCR_LANG = 'chi_tra+eng'
TARGET_DPI = 300
BLACK_THRESHOLD = 100 #稍微放寬黑色的標準，避免深灰字被漏掉

# --- 核心功能 ---

def get_large_image_mask(image_np, text_boxes):
    """產生「圖片禁區遮罩」"""
    img_h, img_w, _ = image_np.shape
    gray = cv2.cvtColor(image_np, cv2.COLOR_BGR2GRAY)
    _, binary = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)
    
    # 把文字塗黑
    for (tx, ty, tw, th) in text_boxes:
        cv2.rectangle(binary, (max(0, tx-5), max(0, ty-5)), (tx+tw+5, ty+th+5), 0, -1)
        
    kernel = np.ones((5,5), np.uint8)
    dilated = cv2.dilate(binary, kernel, iterations=2)
    contours, _ = cv2.findContours(dilated, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    
    danger_zone_mask = np.zeros((img_h, img_w), dtype=np.uint8)
    for cnt in contours:
        x, y, w, h = cv2.boundingRect(cnt)
        area = w * h
        # 面積門檻：頁面的 2%
        if area > (img_w * img_h * 0.02):
            cv2.rectangle(danger_zone_mask, (x, y), (x+w, y+h), 255, -1)
            
    buffer_kernel = np.ones((10, 10), np.uint8) # buffer 稍微縮小一點，避免誤傷主文
    danger_zone_mask = cv2.dilate(danger_zone_mask, buffer_kernel, iterations=1)
    return danger_zone_mask

def is_touching_image(x, y, w, h, danger_mask):
    roi = danger_mask[y:y+h, x:x+w]
    return cv2.countNonZero(roi) > 0

def is_list_item(text):
    text = text.strip()
    if not text: return False
    markers = ['•', '●', '○', '▪', '▫', '◆', '◇', '➢', '➣', '➤', '→', '-', '—', '–', '*', '>']
    if any(text.startswith(m) for m in markers): return True
    pattern = r'^(\d+|[a-zA-Z])[\.\)]\s+'
    if re.match(pattern, text): return True
    return False

def is_text_black(image_np, x, y, w, h):
    img_h, img_w, _ = image_np.shape
    x = max(0, x); y = max(0, y)
    w = min(w, img_w - x); h = min(h, img_h - y)
    if w <= 0 or h <= 0: return False

    roi = image_np[y:y+h, x:x+w]
    gray = cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)
    _, mask = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)
    if cv2.countNonZero(mask) == 0: return False

    mean_val = cv2.mean(roi, mask=mask)
    b, g, r = mean_val[0], mean_val[1], mean_val[2]
    if b < BLACK_THRESHOLD and g < BLACK_THRESHOLD and r < BLACK_THRESHOLD:
        return True
    return False

def get_smart_median_color(image_np, x, y, w, h):
    img_h, img_w, _ = image_np.shape
    sample_w = 10
    x1 = max(0, x - sample_w); x2 = x
    y1 = y; y2 = min(img_h, y + min(h, 10))
    if (x2 - x1) < 2:
        x1 = x; x2 = min(img_w, x + sample_w)
        y1 = max(0, y - 5); y2 = y
    try:
        roi = image_np[y1:y2, x1:x2]
        if roi.size == 0: return (255, 255, 255)
        median_color = np.median(roi, axis=(0, 1))
        return (int(median_color[0]), int(median_color[1]), int(median_color[2]))
    except:
        return (255, 255, 255)

def get_font_size_float(heights_px):
    if not heights_px: return 12.0
    avg_h_px = np.mean(heights_px)
    size_pt = (avg_h_px / TARGET_DPI) * 72 * 0.85
    if size_pt < 9: size_pt = 10
    if size_pt > 120: size_pt = 120
    return size_pt

def process_pdf(uploaded_file):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    bytes_data = uploaded_file.getvalue()
    status_text = st.empty()
    progress_bar = st.progress(0)
    
    status_text.text("正在轉檔與分析 (300 DPI)...")
    images = convert_from_bytes(bytes_data, dpi=TARGET_DPI)
    total_pages = len(images)
    
    for i, img in enumerate(images):
        status_text.text(f"🔄 正在處理第 {i+1} / {total_pages} 頁...")
        
        img_np = np.array(img)
        img_np = cv2.cvtColor(img_np, cv2.COLOR_RGB2BGR)
        img_h, img_w, _ = img_np.shape
        
        data = pytesseract.image_to_data(img, lang=OCR_LANG, output_type=Output.DICT)
        
        paragraphs = {}
        all_text_boxes = [] 
        n_boxes = len(data['text'])
        
        clean_bg_img = img_np.copy()
        
        # 1. 資料收集
        for j in range(n_boxes):
            conf = int(data['conf'][j])
            text = data['text'][j].strip()
            
            if conf > 30 and len(text) > 0:
                x, y, w, h = data['left'][j], data['top'][j], data['width'][j], data['height'][j]
                all_text_boxes.append((x, y, w, h))
                
                key = (data['block_num'][j], data['par_num'][j])
                if key not in paragraphs:
                    paragraphs[key] = {
                        'text_list': [], 'rects': [], 'heights': [], 'line_nums': set()
                    }
                paragraphs[key]['text_list'].append(text)
                paragraphs[key]['rects'].append((x, y, w, h))
                paragraphs[key]['heights'].append(h)
                paragraphs[key]['line_nums'].add(data['line_num'][j])

        # 2. 禁區遮罩與最大字體
        danger_mask = get_large_image_mask(img_np, all_text_boxes)
        
        max_font_size_on_page = 0
        for key in paragraphs:
            f_size = get_font_size_float(paragraphs[key]['heights'])
            paragraphs[key]['calculated_size'] = f_size
            if f_size > max_font_size_on_page:
                max_font_size_on_page = f_size

        # 3. 決策邏輯
        for key, p_data in paragraphs.items():
            full_text = " ".join(p_data['text_list'])
            all_rects = p_data['rects']
            
            min_x = min([r[0] for r in all_rects])
            min_y = min([r[1] for r in all_rects])
            max_x2 = max([r[0] + r[2] for r in all_rects])
            max_y2 = max([r[1] + r[3] for r in all_rects])
            p_w = max_x2 - min_x
            p_h = max_y2 - min_y
            
            # NotebookLM 移除
            if "notebook" in full_text.lower() and min_y > (img_h * 0.8):
                bg_color = get_smart_median_color(img_np, min_x, min_y, p_w, p_h)
                cv2.rectangle(clean_bg_img, (min_x-2, min_y-2), (max_x2+2, max_y2+2), bg_color, -1)
                continue 

            # 特徵判斷
            is_bullet = is_list_item(full_text)
            is_black = is_text_black(img_np, min_x, min_y, p_w, p_h)
            is_title = (p_data['calculated_size'] >= max_font_size_on_page - 2) and (max_font_size_on_page > 14)
            
            # --- 段落增強判斷 (修正點) ---
            # 只要超過 2 行 -> 是段落
            is_multiline = len(p_data['line_nums']) >= 2
            # 就算只有 1 行，如果字數夠多 (例如 > 10 個中文字或單字) -> 是段落
            word_count = len(full_text) 
            is_long_text = word_count > 10
            # 就算只有 1 行，如果寬度超過版面的 30% -> 是段落
            is_wide_text = p_w > (img_w * 0.3)
            
            # 總合：是否為「實質內文」
            is_content = is_multiline or is_long_text or is_wide_text

            is_touching_img = is_touching_image(min_x, min_y, p_w, p_h, danger_mask)
            
            should_extract = False
            
            # --- 最終決策樹 ---
            if is_bullet:
                should_extract = True # 清單無敵，必拆
            elif is_title:
                should_extract = True # 標題無敵，必拆
            elif is_black:
                if is_content:
                    # 如果是實質內文 (長句/寬句/多行)，即使碰到圖片也要拆
                    # 因為通常主文排版都會貼著圖片，不能因為稍微碰到就不拆
                    should_extract = True
                else:
                    # 如果是「短、窄、單行」的黑字
                    # 這時候才檢查有沒有碰到圖片
                    # 碰到圖片 -> 圖說 (不拆)
                    # 沒碰到圖片 -> 可能是頁碼或雜訊 (這裡選擇不拆，保持乾淨)
                    if not is_touching_img:
                        pass 
            
            if should_extract:
                bg_color = get_smart_median_color(img_np, min_x, min_y, p_w, p_h)
                cv2.rectangle(clean_bg_img, (min_x-2, min_y-2), (max_x2+2, max_y2+2), bg_color, -1)
                p_data['should_export'] = True
                p_data['bbox'] = (min_x, min_y, p_w, p_h)
            else:
                p_data['should_export'] = False

        # 4. 產生 PPT
        clean_bg_rgb = cv2.cvtColor(clean_bg_img, cv2.COLOR_BGR2RGB)
        pil_img = Image.fromarray(clean_bg_rgb)
        img_stream = io.BytesIO()
        pil_img.save(img_stream, format='JPEG', quality=95)
        img_stream.seek(0)
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        scale_x = prs.slide_width / img_w
        scale_y = prs.slide_height / img_h
        
        for key, p_data in paragraphs.items():
            if not p_data.get('should_export'): continue
                
            min_x, min_y, p_w, p_h = p_data['bbox']
            full_text = " ".join(p_data['text_list'])
            
            ppt_x = min_x * scale_x
            ppt_y = min_y * scale_y
            ppt_w = p_w * scale_x + Inches(0.15)
            ppt_h = p_h * scale_y
            this_font_size = p_data['calculated_size']

            try:
                txBox = slide.shapes.add_textbox(ppt_x, ppt_y, ppt_w, ppt_h)
                tf = txBox.text_frame
                tf.word_wrap = True
                tf.text = full_text
                for paragraph in tf.paragraphs:
                    paragraph.font.size = Pt(this_font_size)
                    paragraph.font.name = "Arial"
                    paragraph.font.color.rgb = RGBColor(0, 0, 0)
                    if (this_font_size >= max_font_size_on_page - 2) and (max_font_size_on_page > 14):
                        paragraph.font.bold = True
                    else:
                        paragraph.font.bold = False
            except:
                pass
        
        progress_bar.progress((i + 1) / total_pages)

    status_text.text("✅ 轉換完成！")
    ppt_output = io.BytesIO()
    prs.save(ppt_output)
    ppt_output.seek(0)
    return ppt_output

# --- 介面 ---
uploaded_file = st.file_uploader("📂 請上傳 PDF 檔案", type=["pdf"])

if uploaded_file is not None:
    if st.button("🚀 開始轉換"):
        try:
            original_filename = uploaded_file.name
            file_root, _ = os.path.splitext(original_filename)
            new_filename = f"{file_root}_Fixed.pptx"

            ppt_file = process_pdf(uploaded_file)
            st.success(f"🎉 處理成功！")
            
            st.download_button(
                label=f"📥 下載 {new_filename}",
                data=ppt_file,
                file_name=new_filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        except Exception as e:
            st.error(f"❌ 錯誤：{e}")
            st.info("💡 提示：如果線上報錯，請檢查 requirements.txt 是否包含 opencv-python-headless。")
